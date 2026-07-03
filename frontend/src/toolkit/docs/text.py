# /// script
# requires-python = ">=3.10"
# dependencies = ["python-docx>=1.1", "python-pptx>=0.6", "odfpy>=1.4"]
# ///
"""Extract text from an Office document — .docx (paragraphs), .pptx (per-slide),
.odt (paragraphs). One tool, three formats.

Usage: uv run text.py <document> [--json]
Output: {format, text, n_paragraphs|n_slides, paragraphs|slides:[...]}
"""
import argparse, json, os, sys


def from_docx(path):
    import docx
    d = docx.Document(path)
    paras = [p.text.strip() for p in d.paragraphs if p.text.strip()]
    return {"format": "docx", "n_paragraphs": len(paras), "paragraphs": paras,
            "text": "\n".join(paras)}


def from_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        chunks = [sh.text.strip() for sh in slide.shapes if sh.has_text_frame and sh.text.strip()]
        slides.append({"slide": i + 1, "text": "\n".join(chunks)})
    return {"format": "pptx", "n_slides": len(slides), "slides": slides,
            "text": "\n\n".join(f"[slide {s['slide']}]\n{s['text']}" for s in slides)}


def from_odt(path):
    from odf.opendocument import load
    from odf import text as odftext
    from odf.element import Element
    doc = load(path)
    paras = []
    for p in doc.getElementsByType(odftext.P):
        t = "".join(n.data for n in p.childNodes if n.nodeType == 3 or hasattr(n, "data"))
        if t.strip():
            paras.append(t.strip())
    return {"format": "odt", "n_paragraphs": len(paras), "paragraphs": paras,
            "text": "\n".join(paras)}


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("document")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()
    ext = os.path.splitext(args.document)[1].lower()
    try:
        if ext == ".docx":
            result = from_docx(args.document)
        elif ext == ".pptx":
            result = from_pptx(args.document)
        elif ext == ".odt":
            result = from_odt(args.document)
        else:
            result = {"error": f"unsupported docs extension {ext}; try docs/structure.py or run_python"}
    except Exception as e:
        print(json.dumps({"error": str(e), "type": type(e).__name__})); sys.exit(1)
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
